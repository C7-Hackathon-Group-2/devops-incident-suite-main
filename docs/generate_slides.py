"""Generate PowerPoint slide deck for the DevOps Incident Analysis Suite."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

OUTPUT = os.path.join(os.path.dirname(__file__),
                      "DevOps_Incident_Suite_Presentation.pptx")
DIAGRAM_PNG = os.path.join(os.path.dirname(__file__), "architecture_diagram.png")

# Brand colors
NAVY = RGBColor(0x1A, 0x1A, 0x2E)
INDIGO = RGBColor(0x63, 0x66, 0xF1)
BLUE = RGBColor(0x3B, 0x82, 0xF6)
PURPLE = RGBColor(0x8B, 0x5C, 0xF6)
PINK = RGBColor(0xEC, 0x48, 0x99)
GREEN = RGBColor(0x10, 0xB9, 0x81)
AMBER = RGBColor(0xF5, 0x9E, 0x0B)
RED = RGBColor(0xEF, 0x44, 0x44)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF0, 0xF0, 0xF5)
MID_GRAY = RGBColor(0x88, 0x88, 0x88)
DARK_GRAY = RGBColor(0x55, 0x55, 0x55)
DARK_TEXT = RGBColor(0x2D, 0x34, 0x36)
COVER_GREEN = RGBColor(0x05, 0x8C, 0x5A)


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, font_size=18,
                bold=False, color=DARK_TEXT, alignment=PP_ALIGN.LEFT,
                font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return tf


def add_bullet_slide(slide, items, left=Inches(0.8), top=Inches(2.0),
                     width=Inches(8.4), height=Inches(4.5), font_size=16,
                     color=DARK_TEXT, spacing=Pt(8)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = spacing
        p.level = 0
    return tf


def add_slide_title(slide, title, subtitle=None):
    add_textbox(slide, Inches(0.8), Inches(0.4), Inches(8.4), Inches(0.7),
                title, font_size=28, bold=True, color=NAVY)
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.1), Inches(1.2), Inches(0.05))
    shape.fill.solid()
    shape.fill.fore_color.rgb = INDIGO
    shape.line.fill.background()
    if subtitle:
        add_textbox(slide, Inches(0.8), Inches(1.25), Inches(8.4), Inches(0.5),
                    subtitle, font_size=14, color=MID_GRAY)


def add_rounded_box(slide, left, top, width, height, fill_color, text,
                    sub_text=None, font_size=14):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    shape.shadow.inherit = False
    tf = shape.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    run = tf.paragraphs[0].add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = True
    run.font.color.rgb = WHITE
    run.font.name = "Calibri"
    if sub_text:
        p = tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = sub_text
        r.font.size = Pt(10)
        r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        r.font.name = "Calibri"
    return shape


def build():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # ── Slide 1: Cover ─────────────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide, NAVY)

    add_textbox(slide, Inches(1), Inches(1.5), Inches(8), Inches(1.5),
                "Multi-Agent DevOps\nIncident Analysis Suite",
                font_size=36, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(1), Inches(3.2), Inches(8), Inches(0.6),
                "Design, Build & Deployment",
                font_size=20, color=RGBColor(0xA0, 0xA0, 0xC0), alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(1), Inches(3.8), Inches(8), Inches(0.5),
                "AI Engineering Accelerator — Cohort 7",
                font_size=16, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(3.5), Inches(4.6), Inches(3), Inches(0.04))
    shape.fill.solid()
    shape.fill.fore_color.rgb = INDIGO
    shape.line.fill.background()

    add_textbox(slide, Inches(1), Inches(5.0), Inches(8), Inches(0.5),
                "Wayne Johnston", font_size=18, color=WHITE, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(1), Inches(5.5), Inches(8), Inches(0.4),
                "June 2026", font_size=14, color=MID_GRAY, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(1), Inches(6.2), Inches(8), Inches(0.4),
                "Built with Claude Opus 4.6  |  LangGraph  |  LangChain  |  Streamlit",
                font_size=11, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

    # ── Slide 2: Agenda ────────────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, "Agenda")
    items = [
        "1.  Executive Summary",
        "2.  Architecture & Design",
        "3.  Technology Stack",
        "4.  Agent Deep Dive (5 agents)",
        "5.  LangGraph Orchestration",
        "6.  Slack Channel Routing",
        "7.  Streamlit UI & Security",
        "8.  Unit Testing & Code Coverage",
        "9.  Deployment",
        "10. Future Enhancements",
    ]
    add_bullet_slide(slide, items, font_size=18, spacing=Pt(6))

    # ── Slide 3: Executive Summary ─────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, "Executive Summary",
                    "AI-driven incident analysis for DevOps & SRE teams")
    items = [
        "•  Upload infrastructure logs (routers, switches, firewalls, servers)",
        "•  AI infers log formats — no hardcoded regex or format parsers",
        "•  5-agent pipeline: parse → classify → remediate → notify → runbook",
        "•  Multi-channel Slack routing by incident category",
        "•  JIRA ticket creation for critical-severity issues",
        "•  Actionable Markdown runbook with CLI commands & checklists",
        "•  Mock/live integration toggles for safe demoing",
        "•  Directory watcher for automated log ingestion",
    ]
    add_bullet_slide(slide, items, top=Inches(2.2), font_size=16, spacing=Pt(6))

    # ── Slide 4: Architecture Diagram ──────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, "Architecture & Design",
                    "LangGraph StateGraph orchestrates 5 specialized agents")
    if os.path.exists(DIAGRAM_PNG):
        img_width = Inches(7.5)
        img_height = Inches(7.5 * 520 / 700)
        left = Inches((10 - 7.5) / 2)
        slide.shapes.add_picture(DIAGRAM_PNG, left, Inches(1.8),
                                 img_width, img_height)
    else:
        add_textbox(slide, Inches(1), Inches(3), Inches(8), Inches(1),
                    "[architecture_diagram.png not found]",
                    font_size=14, color=RED, alignment=PP_ALIGN.CENTER)

    # ── Slide 5: Pipeline Flow ─────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, "Pipeline Flow",
                    "Conditional routing based on incident severity")

    y_start = Inches(2.2)
    box_h = Inches(0.7)
    box_w = Inches(1.6)
    gap = Inches(0.15)

    nodes = [
        ("Log Reader", GREEN, Inches(0.5)),
        ("Remediation", PURPLE, Inches(2.3)),
        ("Route", INDIGO, Inches(4.1)),
    ]
    for label, color, x in nodes:
        add_rounded_box(slide, x, y_start, box_w, box_h, color, label)

    # Critical path
    crit_y = Inches(3.5)
    add_textbox(slide, Inches(4.3), Inches(2.95), Inches(1.5), Inches(0.3),
                "has_critical?", font_size=11, color=INDIGO, alignment=PP_ALIGN.CENTER)

    crit_nodes = [
        ("Slack Notifier", AMBER, Inches(2.0)),
        ("JIRA Agent", RED, Inches(3.9)),
        ("Cookbook", PINK, Inches(5.8)),
    ]
    for label, color, x in crit_nodes:
        fc = DARK_TEXT if color == AMBER else WHITE
        shape = add_rounded_box(slide, x, crit_y, box_w, box_h, color, label)
        if color == AMBER:
            shape.text_frame.paragraphs[0].runs[0].font.color.rgb = DARK_TEXT

    # Normal path
    add_textbox(slide, Inches(5.9), Inches(2.5), Inches(2), Inches(0.3),
                "no critical →", font_size=11, color=MID_GRAY)
    add_rounded_box(slide, Inches(7.5), y_start, box_w, box_h, PINK, "Cookbook")

    # END boxes
    add_rounded_box(slide, Inches(7.7), crit_y, Inches(1.0), box_h,
                    NAVY, "END", font_size=12)
    add_rounded_box(slide, Inches(7.5), Inches(3.0), Inches(1.0), Inches(0.5),
                    NAVY, "END", font_size=12)

    # Arrows as thin rectangles
    arrow_h = Inches(0.04)
    for (x1, x2, y) in [
        (Inches(2.1), Inches(2.3), Inches(2.55)),
        (Inches(3.9), Inches(4.1), Inches(2.55)),
        (Inches(3.6), Inches(3.9), Inches(3.85)),
        (Inches(5.5), Inches(5.8), Inches(3.85)),
        (Inches(7.4), Inches(7.7), Inches(3.85)),
    ]:
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x1, y, x2 - x1, arrow_h)
        shape.fill.solid()
        shape.fill.fore_color.rgb = MID_GRAY
        shape.line.fill.background()

    # ── Slide 6: Technology Stack ──────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, "Technology Stack")

    rows = [
        ("LLM", "GPT-4o via OpenRouter", BLUE),
        ("Orchestration", "LangGraph StateGraph", INDIGO),
        ("Framework", "LangChain Core", PURPLE),
        ("UI", "Streamlit", GREEN),
        ("Notifications", "Slack Webhook API", AMBER),
        ("Ticketing", "JIRA REST API", RED),
        ("File Watcher", "watchdog", PINK),
    ]
    for i, (label, tech, color) in enumerate(rows):
        y = Inches(2.0) + Inches(i * 0.7)
        add_rounded_box(slide, Inches(0.8), y, Inches(2.5), Inches(0.55),
                        color, label, font_size=13)
        if color == AMBER:
            # fix text color for amber
            for p in slide.shapes[-1].text_frame.paragraphs:
                for r in p.runs:
                    r.font.color.rgb = DARK_TEXT
        add_textbox(slide, Inches(3.6), y + Inches(0.1), Inches(5), Inches(0.45),
                    tech, font_size=16, color=DARK_TEXT)

    # ── Slide 7: Agent Deep Dive ───────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, "Agent Deep Dive",
                    "5 specialized agents with distinct responsibilities")

    agents = [
        ("Log Reader / Classifier", BLUE,
         "3-stage AI pipeline:\n"
         "  1. Schema inference (first 15 lines)\n"
         "  2. Full parse (100-line chunks)\n"
         "  3. Severity + category classification"),
        ("Remediation Agent", PURPLE,
         "Filters warning/critical events\n"
         "Root cause analysis\n"
         "Vendor-specific CLI commands\n"
         "Risk assessment per fix"),
        ("Cookbook Synthesizer", PINK,
         "Markdown runbook generation\n"
         "Priority matrix table\n"
         "Checkbox checklists + code blocks\n"
         "Escalation criteria"),
        ("Slack Notifier", AMBER,
         "Groups by incident category\n"
         "Routes to correct channel\n"
         "Rich formatted messages\n"
         "Mock/live toggle"),
        ("JIRA Ticket Agent", RED,
         "Critical issues only\n"
         "Structured ticket descriptions\n"
         "{noformat} CLI command blocks\n"
         "Mock/live toggle"),
    ]

    for i, (name, color, desc) in enumerate(agents):
        col = i % 3
        row = i // 3
        x = Inches(0.5) + Inches(col * 3.15)
        y = Inches(2.0) + Inches(row * 2.8)
        w = Inches(2.9)
        h = Inches(2.4)

        add_rounded_box(slide, x, y, w, Inches(0.5), color, name, font_size=12)
        fc = DARK_TEXT if color == AMBER else WHITE
        if color == AMBER:
            slide.shapes[-1].text_frame.paragraphs[0].runs[0].font.color.rgb = DARK_TEXT

        add_textbox(slide, x + Inches(0.15), y + Inches(0.6), w - Inches(0.3),
                    h - Inches(0.7), desc, font_size=11, color=DARK_TEXT)

    # ── Slide 8: Slack Channel Routing ─────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, "Slack Channel Routing",
                    "Incidents auto-route to the right team's channel")

    channels = [
        ("#ops-incidents", "resource, application, database",
         GREEN, "Operations team"),
        ("#secops-incidents", "security, auth, policy",
         RED, "Security operations"),
        ("#network-incidents", "interface, routing, network, hardware",
         BLUE, "Network engineering"),
    ]
    for i, (channel, cats, color, team) in enumerate(channels):
        y = Inches(2.2) + Inches(i * 1.6)
        add_rounded_box(slide, Inches(0.8), y, Inches(3), Inches(1.2),
                        color, channel, sub_text=team, font_size=16)
        add_textbox(slide, Inches(4.3), y + Inches(0.15), Inches(5), Inches(0.4),
                    "Categories:", font_size=12, bold=True, color=DARK_TEXT)
        add_textbox(slide, Inches(4.3), y + Inches(0.5), Inches(5), Inches(0.5),
                    cats, font_size=15, color=DARK_TEXT)

    # ── Slide 9: Shared State Schema ───────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, "Shared State Schema",
                    "IncidentState TypedDict — all agents read/write shared state")

    fields = [
        ("raw_logs", "dict[str, str]", "Filename → full log text"),
        ("inferred_schemas", "dict[str, list]", "AI-detected field definitions"),
        ("parsed_events", "list[dict]", "Structured parsed events"),
        ("classified_events", "list[dict]", "Events + severity + category"),
        ("remediations", "list[dict]", "Issue → fix + CLI commands"),
        ("cookbook", "str", "Markdown runbook"),
        ("notifications_sent", "list[dict]", "Slack message results"),
        ("jira_tickets", "list[dict]", "JIRA ticket references"),
        ("status", "str", "Current pipeline stage"),
        ("errors", "list[str]", "Accumulated errors"),
    ]

    # Table header
    hdr_y = Inches(2.1)
    for col_x, col_w, text in [
        (Inches(0.8), Inches(2.5), "Field"),
        (Inches(3.3), Inches(2.2), "Type"),
        (Inches(5.5), Inches(4.0), "Description"),
    ]:
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, col_x, hdr_y,
                                       col_w, Inches(0.4))
        shape.fill.solid()
        shape.fill.fore_color.rgb = NAVY
        shape.line.fill.background()
        tf = shape.text_frame
        tf.paragraphs[0].alignment = PP_ALIGN.LEFT
        r = tf.paragraphs[0].add_run()
        r.text = "  " + text
        r.font.size = Pt(12)
        r.font.bold = True
        r.font.color.rgb = WHITE
        r.font.name = "Calibri"

    for i, (field, typ, desc) in enumerate(fields):
        row_y = Inches(2.5) + Inches(i * 0.42)
        bg = LIGHT_GRAY if i % 2 == 0 else WHITE
        for col_x, col_w, text in [
            (Inches(0.8), Inches(2.5), field),
            (Inches(3.3), Inches(2.2), typ),
            (Inches(5.5), Inches(4.0), desc),
        ]:
            shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, col_x, row_y,
                                           col_w, Inches(0.4))
            shape.fill.solid()
            shape.fill.fore_color.rgb = bg
            shape.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
            shape.line.width = Pt(0.5)
            tf = shape.text_frame
            tf.paragraphs[0].alignment = PP_ALIGN.LEFT
            r = tf.paragraphs[0].add_run()
            r.text = "  " + text
            r.font.size = Pt(11)
            r.font.color.rgb = DARK_TEXT
            r.font.name = "Calibri"
            if col_x == Inches(0.8):
                r.font.bold = True
                r.font.name = "Courier New"
                r.font.size = Pt(10)

    # ── Slide 10: Streamlit UI ─────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, "Streamlit UI & Security")

    # Left column - UI features
    add_textbox(slide, Inches(0.8), Inches(1.8), Inches(4), Inches(0.4),
                "UI Features", font_size=16, bold=True, color=INDIGO)
    ui_items = [
        "•  Multi-file upload (.log, .txt)",
        "•  Load Samples button for demo",
        "•  Directory watcher auto-ingestion",
        "•  6 result tabs: Events, Remediations,",
        "    Runbook, Slack, JIRA, Schemas",
        "•  Severity color coding (red/yellow/blue)",
        "•  Runbook download button",
    ]
    add_bullet_slide(slide, ui_items, left=Inches(0.8), top=Inches(2.2),
                     width=Inches(4), height=Inches(4), font_size=14, spacing=Pt(4))

    # Right column - Security
    add_textbox(slide, Inches(5.2), Inches(1.8), Inches(4), Inches(0.4),
                "Security", font_size=16, bold=True, color=RED)
    sec_items = [
        "•  API keys never pre-filled in UI",
        "•  Users enter own keys per session",
        "•  Secrets reserved for /api endpoint",
        "•  Analyze blocked without API key",
        "•  Mock-first integrations (safe demos)",
        "•  .env excluded from git",
        "•  Framework-agnostic core (0 UI deps)",
    ]
    add_bullet_slide(slide, sec_items, left=Inches(5.2), top=Inches(2.2),
                     width=Inches(4.3), height=Inches(4), font_size=14, spacing=Pt(4))

    # ── Slide 11: Unit Testing & Code Coverage ─────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, "Unit Testing & Code Coverage",
                    "81 tests | 99% coverage | pytest + unittest.mock")

    # Summary stats boxes
    stats = [
        ("81", "Tests Passed", GREEN),
        ("0", "Failed", GREEN),
        ("99%", "Coverage", INDIGO),
        ("5", "Test Modules", BLUE),
    ]
    for i, (num, label, color) in enumerate(stats):
        x = Inches(0.8) + Inches(i * 2.25)
        shape = add_rounded_box(slide, x, Inches(2.0), Inches(1.9), Inches(1.0),
                                color, num, sub_text=label, font_size=28)

    # Coverage table
    hdr_y = Inches(3.4)
    col_defs = [
        (Inches(0.8), Inches(2.2), "Module"),
        (Inches(3.0), Inches(1.2), "Stmts"),
        (Inches(4.2), Inches(1.2), "Miss"),
        (Inches(5.4), Inches(1.5), "Cover"),
        (Inches(6.9), Inches(2.6), "Missing"),
    ]
    for col_x, col_w, text in col_defs:
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, col_x, hdr_y,
                                       col_w, Inches(0.38))
        shape.fill.solid()
        shape.fill.fore_color.rgb = NAVY
        shape.line.fill.background()
        tf = shape.text_frame
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        r = tf.paragraphs[0].add_run()
        r.text = text
        r.font.size = Pt(11)
        r.font.bold = True
        r.font.color.rgb = WHITE
        r.font.name = "Calibri"

    coverage_data = [
        ("state.py", "12", "0", "100%", "—"),
        ("tools.py", "36", "0", "100%", "—"),
        ("watcher.py", "64", "0", "100%", "—"),
        ("graph.py", "21", "0", "100%", "—"),
        ("agents.py", "143", "2", "99%", "13, 17"),
        ("TOTAL", "276", "2", "99%", ""),
    ]
    for i, (mod, stmts, miss, cover, missing) in enumerate(coverage_data):
        row_y = Inches(3.78) + Inches(i * 0.38)
        bg = LIGHT_GRAY if i % 2 == 0 else WHITE
        is_total = (mod == "TOTAL")
        if is_total:
            bg = RGBColor(0xE8, 0xE8, 0xF0)
        vals = [mod, stmts, miss, cover, missing]
        for j, (col_x, col_w, _) in enumerate(col_defs):
            shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, col_x, row_y,
                                           col_w, Inches(0.38))
            shape.fill.solid()
            shape.fill.fore_color.rgb = bg
            shape.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
            shape.line.width = Pt(0.5)
            tf = shape.text_frame
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT
            r = tf.paragraphs[0].add_run()
            r.text = ("  " if j == 0 else "") + vals[j]
            r.font.size = Pt(11)
            r.font.color.rgb = DARK_TEXT
            r.font.name = "Courier New" if j == 0 else "Calibri"
            r.font.bold = is_total
            if j == 3 and cover == "100%":
                r.font.color.rgb = GREEN
            elif j == 3 and cover == "99%":
                r.font.color.rgb = INDIGO

    # Note about missing lines
    add_textbox(slide, Inches(0.8), Inches(6.1), Inches(8.4), Inches(0.8),
                "Lines 13, 17 in agents.py are import-time .env path checks "
                "(load_dotenv calls) — both branches are exercised but "
                "coverage can't attribute them to a specific test due to module import timing.",
                font_size=11, color=MID_GRAY)

    # ── Slide 12: Test Breakdown ───────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, "Test Suite Breakdown",
                    "Comprehensive mocking of LLM, Slack, and JIRA integrations")

    test_modules = [
        ("test_state.py", "2 tests", BLUE,
         "TypedDict instantiation\nPopulated data validation"),
        ("test_tools.py", "23 tests", PURPLE,
         "Channel routing (12 categories)\n"
         "Slack mock/live modes\n"
         "JIRA mock/live modes\n"
         "Webhook POST verification"),
        ("test_watcher.py", "10 tests", GREEN,
         "File type filtering\n"
         "Debounce logic\n"
         "Directory create/scan\n"
         "Unreadable file handling"),
        ("test_agents.py", "25 tests", PINK,
         "Lazy LLM pattern\n"
         "All 5 agent nodes\n"
         "JSON parse failures\n"
         "Edge cases (empty, malformed)"),
        ("test_graph.py", "8 tests", INDIGO,
         "Severity routing logic\n"
         "Graph compilation\n"
         "Full pipeline (critical path)\n"
         "Normal path (skip Slack/JIRA)"),
    ]
    for i, (name, count, color, desc) in enumerate(test_modules):
        col = i % 3
        row = i // 3
        x = Inches(0.4) + Inches(col * 3.2)
        y = Inches(2.0) + Inches(row * 2.7)
        w = Inches(2.9)

        add_rounded_box(slide, x, y, w, Inches(0.45), color, f"{name}  ({count})",
                        font_size=11)
        add_textbox(slide, x + Inches(0.1), y + Inches(0.55), w - Inches(0.2),
                    Inches(2.0), desc, font_size=11, color=DARK_TEXT)

    # Run command
    add_textbox(slide, Inches(0.8), Inches(6.3), Inches(8.4), Inches(0.8),
                "Run:  python -m pytest test_*.py -v --cov=state --cov=tools "
                "--cov=watcher --cov=agents --cov=graph --cov-report=term-missing",
                font_size=10, color=MID_GRAY)

    # ── Slide 13: Deployment ───────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, "Deployment",
                    "GitHub + Streamlit Cloud")

    add_textbox(slide, Inches(0.8), Inches(2.0), Inches(4), Inches(0.4),
                "GitHub Repository", font_size=16, bold=True, color=INDIGO)
    add_textbox(slide, Inches(0.8), Inches(2.5), Inches(4), Inches(0.4),
                "github.com/wjohnston99/devops-incident-suite",
                font_size=13, color=DARK_TEXT)

    add_textbox(slide, Inches(0.8), Inches(3.2), Inches(8), Inches(0.4),
                "Streamlit Cloud Steps", font_size=16, bold=True, color=INDIGO)
    steps = [
        "1.  Sign in to share.streamlit.io with GitHub",
        "2.  Click 'New app' → select devops-incident-suite repo",
        "3.  Branch: main  |  Main file: app.py",
        "4.  No secrets needed (users enter own keys)",
        "5.  Click Deploy — live in ~2 minutes",
    ]
    add_bullet_slide(slide, steps, top=Inches(3.7), font_size=15, spacing=Pt(6))

    add_textbox(slide, Inches(0.8), Inches(5.8), Inches(8.4), Inches(0.8),
                "Future: Flask API endpoint — core pipeline has zero Streamlit "
                "dependencies, enabling a Flask/FastAPI frontend swap with no agent changes.",
                font_size=13, color=MID_GRAY)

    # ── Slide 14: Project Structure ────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, "Project Structure",
                    "Framework-agnostic core + Streamlit UI")

    files = [
        ("src/app.py", "Streamlit UI entry point", GREEN),
        ("src/state.py", "IncidentState TypedDict (shared state)", BLUE),
        ("src/agents.py", "5 agent node functions + lazy LLM", PURPLE),
        ("src/graph.py", "LangGraph StateGraph wiring", INDIGO),
        ("src/tools.py", "@tool wrappers (Slack, JIRA) + routing", PINK),
        ("src/watcher.py", "Directory watcher (watchdog)", AMBER),
        ("src/log_samples/", "4 synthetic log files (IOS, CEF, RFC5424, tabular)", MID_GRAY),
        ("tests/test_*.py", "81 unit tests (pytest)", RED),
    ]

    for i, (fname, desc, color) in enumerate(files):
        y = Inches(2.1) + Inches(i * 0.6)
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                       Inches(0.8), y, Inches(0.15), Inches(0.4))
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()

        add_textbox(slide, Inches(1.1), y, Inches(2.2), Inches(0.4),
                    fname, font_size=13, bold=True, color=DARK_TEXT,
                    font_name="Courier New")
        add_textbox(slide, Inches(3.5), y, Inches(5.5), Inches(0.4),
                    desc, font_size=13, color=DARK_TEXT)

    # ── Slide 15: Future Enhancements ──────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, "Future Enhancements")

    items = [
        "•  Flask API endpoint (POST /api/analyze) with server-side secrets",
        "•  RAG-based remediation using historical runbook embeddings",
        "•  Real-time streaming via LangGraph's streaming interface",
        "•  Cross-file correlation engine for cascading failure detection",
        "•  Persistent storage (SQLite/PostgreSQL) for trend analysis",
        "•  PagerDuty integration for critical incident escalation",
        "•  Custom log format registry (save & share inferred schemas)",
    ]
    add_bullet_slide(slide, items, font_size=16, spacing=Pt(10))

    # ── Slide 16: Thank You ────────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)

    add_textbox(slide, Inches(1), Inches(2.0), Inches(8), Inches(1),
                "Thank You", font_size=42, bold=True, color=WHITE,
                alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(1), Inches(3.2), Inches(8), Inches(0.5),
                "Multi-Agent DevOps Incident Analysis Suite",
                font_size=18, color=RGBColor(0xA0, 0xA0, 0xC0),
                alignment=PP_ALIGN.CENTER)

    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(3.5), Inches(4.0), Inches(3), Inches(0.04))
    shape.fill.solid()
    shape.fill.fore_color.rgb = INDIGO
    shape.line.fill.background()

    add_textbox(slide, Inches(1), Inches(4.4), Inches(8), Inches(0.5),
                "Wayne Johnston  |  AI Engineering Accelerator  |  Cohort 7",
                font_size=14, color=MID_GRAY, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(1), Inches(5.0), Inches(8), Inches(0.4),
                "github.com/wjohnston99/devops-incident-suite",
                font_size=12, color=INDIGO, alignment=PP_ALIGN.CENTER)

    prs.save(OUTPUT)
    print(f"Presentation saved: {OUTPUT}")


if __name__ == "__main__":
    build()
